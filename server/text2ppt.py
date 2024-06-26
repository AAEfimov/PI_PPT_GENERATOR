"""
text2ppt.py
generate ppt
"""

__author__ = "UrFU team"
__copyright__ = "Copyright 2024, Planet Earth"

import io
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt

import addphoto

default_font = {"name": "Arial", "size": 12, "bold": False, "italic": False}
tempalte_default = "default.pptx"


def find_template(name):
    for root, dirs, files in os.walk("./"):
        if name in files:
            return os.path.join(root, name)


# Create a new PowerPoint presentation
def presentate(
    defined_list,
    img=None,
    title="",
    subtitle="",
    layout=None,
    font_param=default_font,
    templ_name=tempalte_default,
):

    template_path = find_template(templ_name)
    prs = Presentation(template_path)

    # print(defined_list)

    def add_slide(prs, layout, title, subtitle):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = subtitle
        font = slide.shapes.title.text_frame.paragraphs[0].font
        font.name = font_param["name"]
        font.size = Pt(font_param["size"])
        font.bold = font_param["bold"]
        font.italic = font_param["italic"]
        for x in slide.placeholders[1].text_frame.paragraphs:
            font1 = x.font
            font1.name = font_param["name"]
            font1.size = Pt(font_param["size"])
            font1.bold = font_param["bold"]
            font1.italic = font_param["italic"]
        return slide

    def add_slide_img(prs, layout, img_path):
        slide = prs.slides.add_slide(layout)
        img_path = "" + img_path
        left = Inches(1.10)
        top = Inches(0.7)
        width = Inches(8)
        height = Inches(6)
        slide.shapes.add_picture(img_path, left, top, width, height)

    title_slide_layout = prs.slide_layouts[1]  # Only text
    title_slide_layimg = prs.slide_layouts[6]  # Only image

    for d in defined_list:
        d["Summary"] = [
            re.sub(r"\d+\.\s+", "", item).strip()
            for item in d["Summary"]
            if re.sub(r"\d+\.\s+", "", item).strip()
        ]

    for i in range(0, len(defined_list)):
        add_slide(
            prs,
            title_slide_layout,
            defined_list[i]["Topic"] if len(title) == 0 else title,
            "\n".join(
                defined_list[i]["Summary"][0:len(defined_list[i]["Summary"]) // 2]
            ),
        )

        add_slide(
            prs,
            title_slide_layout,
            defined_list[i]["Topic"] if len(subtitle) == 0 else subtitle,
            "\n".join(
                defined_list[i]["Summary"][len(defined_list[i]["Summary"]) // 2:]
            ),
        )

        try:
            imgout = img if img else addphoto.get_images(defined_list[i]["Topic"], 1)[0]
            add_slide_img(prs, title_slide_layimg, "images/" + imgout)
        except Exception as e:
            print("got Exception ", e)
            imgout = img if img else addphoto.get_images(defined_list[i]["Topic"], 1)[0]
            add_slide_img(
                prs,
                title_slide_layimg,
                "images/" + addphoto.get_images(defined_list[i]["Topic"], 2)[1],
            )
        addphoto.empty_images()

    # Save the presentation
    binary_output = io.BytesIO()

    prs.save(binary_output)

    return binary_output
